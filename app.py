import os
from dotenv import load_dotenv
if not os.getenv("VERCEL") and not os.getenv("RENDER"):
    load_dotenv(override=False)
    load_dotenv(".env.local", override=False)

from flask import Flask, request, render_template, send_file, jsonify
from io import BytesIO
from PyPDF2 import PdfReader
from fpdf import FPDF
import unicodedata
from google.oauth2 import service_account
from googleapiclient.discovery import build
import base64
import re
import tempfile
import time
import uuid
from functools import wraps
from requests import RequestException
import requests
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
import json

from database.sqlite_db import DatabaseManager
from services.embeddingService import EmbeddingService
from services.evaluationService import EvaluationService
from services.questionGenerator import QuestionGenerator
from services.retrievalService import RetrievalService
from vectorstore.vector_store import VectorStore

app = Flask(__name__)

embedding_service = EmbeddingService()
evaluation_service = EvaluationService(embedding_service)
vector_store = VectorStore(embedding_service)
retrieval_service = RetrievalService(vector_store)
DEFAULT_DB_PATH = os.path.join(tempfile.gettempdir(), "queryfy.db") if os.getenv("VERCEL") else "queryfy.db"
db = DatabaseManager(os.getenv("QUERYFY_DB_PATH", DEFAULT_DB_PATH))
LAST_GENERATED_RESULTS = []

# Rate limiting
request_count = {}
RATE_LIMIT = 50  # requests per hour per IP

def rate_limit(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        ip = request.remote_addr
        current_time = time.time()
        
        # Clean old entries
        if ip in request_count:
            request_count[ip] = [t for t in request_count[ip] if current_time - t < 3600]
        else:
            request_count[ip] = []
        
        # Check limit
        if len(request_count[ip]) >= RATE_LIMIT:
            return "Rate limit exceeded. Please try again in an hour.", 429
        
        request_count[ip].append(current_time)
        return f(*args, **kwargs)
    return decorated_function

# Configure AI API
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
GROQ_MODEL = os.getenv("GROQ_MODEL", "llama-3.3-70b-versatile")
GROQ_TEMPERATURE = float(os.getenv("GROQ_TEMPERATURE", "0.7"))
GROQ_MAX_TOKENS = int(os.getenv("GROQ_MAX_TOKENS", "8192"))
GROQ_TIMEOUT = int(os.getenv("GROQ_TIMEOUT", "60"))
GROQ_MAX_RETRIES = int(os.getenv("GROQ_MAX_RETRIES", "3"))
STARTUP_CHECKS = {
    "groq_api_key": "configured" if GROQ_API_KEY else "missing",
    "google_credentials": "pending",
    "vector_store": getattr(vector_store, "status", "unknown"),
    "database": f"sqlite:{db.db_path}",
}

# Google Form and Drive API scopes
SCOPES_FORMS = ['https://www.googleapis.com/auth/forms.body']
SCOPES_DRIVE = ['https://www.googleapis.com/auth/drive']


SERVICE_ACCOUNT_JSON = os.getenv("SERVICE_ACCOUNT_JSON")
SERVICE_ACCOUNT_JSON_BASE64 = os.getenv("SERVICE_ACCOUNT_JSON_BASE64")
SERVICE_ACCOUNT_FILE = os.getenv("SERVICE_ACCOUNT_FILE")
form_service = None
drive_service = None
GOOGLE_FORMS_STATUS = "not configured"

def load_google_service_info():
    """Load Google service-account credentials from env JSON, base64, or a local file."""
    if SERVICE_ACCOUNT_JSON:
        return json.loads(SERVICE_ACCOUNT_JSON), "SERVICE_ACCOUNT_JSON"
    if SERVICE_ACCOUNT_JSON_BASE64:
        decoded = base64.b64decode(SERVICE_ACCOUNT_JSON_BASE64).decode("utf-8")
        return json.loads(decoded), "SERVICE_ACCOUNT_JSON_BASE64"
    if SERVICE_ACCOUNT_FILE and os.path.exists(SERVICE_ACCOUNT_FILE):
        with open(SERVICE_ACCOUNT_FILE, "r", encoding="utf-8") as fh:
            return json.load(fh), "SERVICE_ACCOUNT_FILE"
    return None, "missing SERVICE_ACCOUNT_JSON, SERVICE_ACCOUNT_JSON_BASE64, or SERVICE_ACCOUNT_FILE"

# Authenticate Google services
try:
    service_info, google_forms_source = load_google_service_info()
    if service_info:
        creds_forms = service_account.Credentials.from_service_account_info(service_info, scopes=SCOPES_FORMS)
        creds_drive = service_account.Credentials.from_service_account_info(service_info, scopes=SCOPES_DRIVE)
        form_service = build('forms', 'v1', credentials=creds_forms)
        drive_service = build('drive', 'v3', credentials=creds_drive)
        GOOGLE_FORMS_ENABLED = True
        GOOGLE_FORMS_STATUS = f"configured via {google_forms_source}"
        STARTUP_CHECKS["google_credentials"] = GOOGLE_FORMS_STATUS
    else:
        print(f"Google Forms integration disabled: {google_forms_source}")
        GOOGLE_FORMS_ENABLED = False
        GOOGLE_FORMS_STATUS = google_forms_source
        STARTUP_CHECKS["google_credentials"] = GOOGLE_FORMS_STATUS
    
except Exception as e:
    print(f"Google Forms integration disabled: {e}")
    GOOGLE_FORMS_ENABLED = False
    GOOGLE_FORMS_STATUS = str(e)
    STARTUP_CHECKS["google_credentials"] = GOOGLE_FORMS_STATUS

def log_startup_checks():
    print(f"Queryfy startup: Groq configured = {bool(GROQ_API_KEY)}")
    print(f"Queryfy startup: Groq model = {GROQ_MODEL}")
    print(f"Queryfy startup: Google Forms = {GOOGLE_FORMS_STATUS}")
    print(f"Queryfy startup: Vector store = {STARTUP_CHECKS['vector_store']}")
    print(f"Queryfy startup: Database = {STARTUP_CHECKS['database']}")

log_startup_checks()

def validate_email(email):
    """Validate email format"""
    email_pattern = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
    return re.match(email_pattern, email) is not None

def clean_text_for_generation(text):
    """Normalize extracted text before chunking and retrieval."""
    return VectorStore.clean_text(text)

def extract_text_from_pdf(file_path):
    """Extract text from PDF"""
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return clean_text_for_generation(text)
    except Exception as e:
        print(f"Error extracting text: {e}")
        return ""
def extract_text_from_docx(file_path):
    """Extract text from Word document"""
    try:
        doc = Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"
        return clean_text_for_generation(text)
    except Exception as e:
        print(f"Error extracting DOCX text: {e}")
        return ""

def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint"""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return clean_text_for_generation(text)
    except Exception as e:
        print(f"Error extracting PPTX text: {e}")
        return ""

def extract_text_from_file(file_path, filename):
    """Extract text based on file extension"""
    ext = filename.lower().split('.')[-1]
    if ext == 'pdf':
        return extract_text_from_pdf(file_path)
    elif ext == 'docx':
        return extract_text_from_docx(file_path)
    elif ext == 'pptx':
        return extract_text_from_pptx(file_path)
    elif ext == 'txt':
        with open(file_path, 'r', encoding='utf-8') as f:
            return clean_text_for_generation(f.read())
    else:
        return ""
    
def generate_with_groq(prompt, query):
    """Generate using Groq API (super fast!)"""
    if not GROQ_API_KEY:
        raise RuntimeError("Groq is not configured. Set GROQ_API_KEY in Render or Vercel environment variables.")

    payload = {
        "model": GROQ_MODEL,
        "messages": [
            {"role": "system", "content": prompt},
            {"role": "user", "content": query}
        ],
        "temperature": GROQ_TEMPERATURE,
        "max_tokens": GROQ_MAX_TOKENS
    }
    headers = {
        "Authorization": f"Bearer {GROQ_API_KEY}",
        "Content-Type": "application/json"
    }

    last_error = "Groq request failed"
    for attempt in range(1, GROQ_MAX_RETRIES + 1):
        try:
            response = requests.post(
                "https://api.groq.com/openai/v1/chat/completions",
                headers=headers,
                json=payload,
                timeout=GROQ_TIMEOUT
            )

            if response.status_code == 200:
                result = response.json()
                return result['choices'][0]['message']['content']

            if response.status_code in (401, 403):
                raise RuntimeError("Groq authentication failed. Check GROQ_API_KEY in your deployment environment.")

            message = response.text[:500]
            last_error = f"Groq API error {response.status_code}: {message}"
            print(last_error)

            if response.status_code == 429 or response.status_code >= 500:
                retry_after = response.headers.get("retry-after")
                delay = float(retry_after) if retry_after else min(2 ** attempt, 8)
                time.sleep(delay)
                continue

            raise RuntimeError(last_error)
        except RequestException as exc:
            last_error = f"Could not reach Groq API: {exc}"
            print(last_error)
            if attempt < GROQ_MAX_RETRIES:
                time.sleep(min(2 ** attempt, 8))
                continue
            raise RuntimeError("Groq is temporarily unreachable. Please try again in a moment.") from exc

    raise RuntimeError(f"{last_error}. Please try again shortly.")

def llm_generate(prompt):
    """LLM adapter used by the RAG question generator."""
    return generate_with_groq(prompt, "")

def map_question_type(question_type):
    mapping = {
        "mcq": "MCQ",
        "true/false": "True/False",
        "fill-ups": "Fill in the Blanks",
        "short answer": "Short Answer",
        "subjective": "Long Answer",
        "numerical": "Numerical Problems",
    }
    return mapping.get((question_type or "").strip().lower(), "Short Answer")

def convert_questions_to_text(questions):
    """Convert structured RAG output into the existing export parser format."""
    sections = []
    for index, item in enumerate(questions, 1):
        question = (item.get("question") or "").strip()
        if not question:
            continue

        section = f"{index}. {question}"
        for option_index, option in enumerate(item.get("options") or []):
            section += f"\n   {chr(97 + option_index)}) {option}"

        answer = (item.get("answer") or "").strip()
        explanation = (item.get("explanation") or "").strip()
        difficulty = (item.get("difficulty") or "").strip()

        if answer:
            section += f"\n   **Answer: {answer}**"
        if explanation:
            section += f"\n   **Explanation: {explanation}**"
        if difficulty:
            section += f"\n   **Difficulty: {difficulty}**"
        sections.append(section)
    return "\n\n".join(sections)

def build_rag_pipeline(file_text, filename, subject, question_type, difficulty, chunk_size=500, chunk_overlap=80, top_k=6):
    """Upload -> clean -> chunk -> embed -> store -> retrieve relevant context."""
    doc_id = str(uuid.uuid4())
    metadata = {
        "subject": subject,
        "question_type": question_type,
        "difficulty": difficulty,
        "source": filename,
    }
    chunks = vector_store.add_document(
        doc_id=doc_id,
        filename=filename,
        text=file_text,
        metadata=metadata,
        chunk_size=chunk_size,
        overlap=chunk_overlap,
    )
    query_text = f"{subject} {map_question_type(question_type)} {difficulty} assessment concepts"
    retrieved = retrieval_service.retrieve_with_metadata(query_text, top_k=top_k, document_id=doc_id)
    context_chunks = [item["text"] for item in retrieved] or [chunk["text"] for chunk in chunks[:top_k]]
    return doc_id, chunks, retrieved, context_chunks

def generate_questions_with_answers(prompt, question_type, num_questions):
    """Generate questions WITH answer keys using multiple AI providers"""
    
    # Truncate very long prompts
    if len(prompt) > 15000:
        prompt = prompt[:15000] + "..."
    
    system_instruction = """You are an expert educator creating exam questions with answer keys.
    CRITICAL RULES:
    1. NEVER mention "the text", "the passage", "the document" or similar
    2. Write professional, standalone exam questions
    3. ALWAYS provide correct answers
    4. Use clear, academic language"""
    
    if question_type.lower() == "mcq":
        query = f"""Create {num_questions} multiple-choice questions with answer keys.

CONTENT:
{prompt}

FORMAT (STRICTLY FOLLOW):
1. [Question]
   a) [Option]
   b) [Option]
   c) [Option]
   d) [Option]
   **Answer: [letter]**

2. [Question]
   a) [Option]
   b) [Option]
   c) [Option]
   d) [Option]
   **Answer: [letter]**

REQUIREMENTS:
- Direct questions (no "based on text")
- Exactly 4 options per question
- Mark correct answer clearly
- Mix difficulty levels
- Generate EXACTLY {num_questions} questions"""

    elif question_type.lower() == "true/false":
        query = f"""Create {num_questions} true/false questions with answers.

CONTENT:
{prompt}

FORMAT (STRICTLY FOLLOW):
1. [Statement] (True/False)
   **Answer: [True/False]**

2. [Statement] (True/False)
   **Answer: [True/False]**

REQUIREMENTS:
- Clear statements (no source references)
- Mark correct answer
- Balance true and false
- Generate EXACTLY {num_questions} questions"""

    elif question_type.lower() == "fill-ups":
        query = f"""Create {num_questions} fill-in-the-blank questions with answers.

CONTENT:
{prompt}

FORMAT (STRICTLY FOLLOW):
1. [Sentence with ________]
   **Answer: [correct word/phrase]**

2. [Sentence with ________]
   **Answer: [correct word/phrase]**

REQUIREMENTS:
- Use ________ for blank
- No source references
- Provide correct answer
- Generate EXACTLY {num_questions} questions"""

    elif question_type.lower() == "short answer":
        query = f"""Create {num_questions} short answer questions with model answers.

CONTENT:
{prompt}

FORMAT (STRICTLY FOLLOW):
1. [Question]
   **Model Answer: [2-3 sentence answer with key points]**

2. [Question]
   **Model Answer: [2-3 sentence answer with key points]**

REQUIREMENTS:
- Direct questions
- Provide concise model answers
- Generate EXACTLY {num_questions} questions"""

    elif question_type.lower() == "subjective":
        query = f"""Create {num_questions} essay questions with answer guidelines.

CONTENT:
{prompt}

FORMAT (STRICTLY FOLLOW):
1. [Question]
   **Answer Guidelines: [Key points to cover]**

2. [Question]
   **Answer Guidelines: [Key points to cover]**

REQUIREMENTS:
- Analytical questions
- Provide answer guidelines
- Generate EXACTLY {num_questions} questions"""

    else:  # Numerical
        query = f"""Create {num_questions} numerical/problem-solving questions with solutions.

CONTENT:
{prompt}

FORMAT (STRICTLY FOLLOW):
1. [Problem with all data needed]
   **Solution: [Step-by-step solution with final answer]**

2. [Problem with all data needed]
   **Solution: [Step-by-step solution with final answer]**

REQUIREMENTS:
- Complete problem statements
- Provide full solutions
- Generate EXACTLY {num_questions} questions"""

    try:
        questions_text = generate_with_groq(system_instruction, query)
        if not questions_text:
            raise RuntimeError("Groq generation returned an empty response")
        
        # Clean up
        questions_text = re.sub(r'^(Here are|Here\'s|Below are|Sure).*?[:.\n]', '', questions_text, flags=re.IGNORECASE | re.MULTILINE)
        questions_text = re.sub(r'^#.*$', '', questions_text, flags=re.MULTILINE)
        
        return questions_text.strip()
        
    except Exception as e:
        print(f"Error generating questions: {e}")
        raise

def normalize_export_question(item, index):
    """Normalize structured questions for student and teacher exports."""
    question = (item.get("question") or "").strip()
    options = item.get("options") or []
    answer = (item.get("answer") or "").strip()
    explanation = (item.get("explanation") or "").strip()
    source = (item.get("source_excerpt") or item.get("source_chunk") or "").strip()
    difficulty = (item.get("difficulty") or "Medium").strip()
    return {
        "number": index,
        "question": question,
        "options": [str(option).strip() for option in options if str(option).strip()],
        "answer": answer,
        "explanation": explanation,
        "source_reference": source,
        "difficulty": difficulty,
        "relevance_score": float(item.get("relevance_score") or 0),
        "context_coverage": float(item.get("context_coverage") or 0),
        "confidence_score": float(item.get("confidence_score") or 0),
    }

def parse_questions_for_export(raw_questions, generated_items=None):
    """Prefer structured RAG output; fallback to parsing legacy text output."""
    if generated_items:
        return [
            normalize_export_question(item, index)
            for index, item in enumerate(generated_items, 1)
            if (item.get("question") or "").strip()
        ]

    parsed = []
    parts = re.split(r'(?m)^\s*(\d+)\.\s*', raw_questions or "")
    for i in range(1, len(parts), 2):
        number = len(parsed) + 1
        content = parts[i + 1] if i + 1 < len(parts) else ""
        question_text = content
        answer = ""
        explanation = ""
        difficulty = "Medium"
        source = ""

        markers = [
            ("answer", r'\*\*(?:Answer|Solution|Model Answer|Answer Guidelines):\s*(.*?)\*\*'),
            ("explanation", r'\*\*Explanation:\s*(.*?)\*\*'),
            ("source", r'\*\*Source(?: Chunk)?:\s*(.*?)\*\*'),
            ("difficulty", r'\*\*Difficulty:\s*(.*?)\*\*'),
        ]
        for key, pattern in markers:
            match = re.search(pattern, content, flags=re.IGNORECASE | re.DOTALL)
            if match:
                value = match.group(1).strip()
                if key == "answer":
                    answer = value
                elif key == "explanation":
                    explanation = value
                elif key == "source":
                    source = value
                elif key == "difficulty":
                    difficulty = value

        question_text = re.sub(r'\*\*(?:Answer|Solution|Model Answer|Answer Guidelines|Explanation|Source(?: Chunk)?|Difficulty|Quality Scores):.*?(?:\*\*|$)', '', question_text, flags=re.IGNORECASE | re.DOTALL)
        question_text = question_text.strip()
        if question_text:
            parsed.append({
                "number": number,
                "question": question_text,
                "options": [],
                "answer": answer,
                "explanation": explanation,
                "source_reference": source,
                "difficulty": difficulty,
                "relevance_score": 0.0,
                "context_coverage": 0.0,
                "confidence_score": 0.0,
            })
    return parsed

def summarize_assessment(export_questions, subject):
    total = len(export_questions)
    def average(key):
        values = [q.get(key, 0) for q in export_questions if q.get(key, 0)]
        return round(sum(values) / len(values), 3) if values else 0

    difficulty_distribution = {"Easy": 0, "Medium": 0, "Hard": 0}
    for question in export_questions:
        difficulty = (question.get("difficulty") or "Medium").title()
        if difficulty not in difficulty_distribution:
            difficulty_distribution[difficulty] = 0
        difficulty_distribution[difficulty] += 1

    words = re.findall(r"\b[A-Za-z]{5,}\b", " ".join(q["question"] for q in export_questions).lower())
    stop_words = {"which", "what", "where", "when", "about", "explain", "describe", "identify", "compare", "answer", "question", "following"}
    frequencies = {}
    for word in words:
        if word not in stop_words:
            frequencies[word] = frequencies.get(word, 0) + 1
    top_topics = sorted(frequencies.items(), key=lambda item: item[1], reverse=True)[:8]
    topic_summary = ", ".join(word.title() for word, _ in top_topics) or subject or "General"

    return {
        "total_questions": total,
        "average_relevance": average("relevance_score"),
        "average_coverage": average("context_coverage"),
        "average_confidence": average("confidence_score"),
        "difficulty_distribution": difficulty_distribution,
        "topic_coverage_summary": topic_summary,
    }

def pdf_heading(pdf, title):
    pdf.set_font("Arial", 'B', 16)
    pdf.cell(0, 10, title, ln=True, align="C")
    pdf.ln(6)

def pdf_label_value(pdf, label, value):
    pdf.set_font("Arial", 'B', 11)
    pdf.cell(48, 7, remove_special_characters(label), border=0)
    pdf.set_font("Arial", size=11)
    pdf.multi_cell(0, 7, remove_special_characters(str(value)))

def save_questions_to_pdf(raw_questions, subject, marks, include_answers=True, generated_items=None):
    """Save a clean printable PDF with summary and optional teacher sheet."""
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    export_questions = parse_questions_for_export(raw_questions, generated_items)
    summary = summarize_assessment(export_questions, subject)

    pdf.add_page()
    pdf_heading(pdf, "ASSESSMENT PAPER")
    pdf.set_font("Arial", 'B', 12)
    pdf.cell(0, 8, remove_special_characters(f"Subject: {subject}"), ln=True)
    pdf.cell(0, 8, remove_special_characters(f"Maximum Marks: {marks}"), ln=True)
    pdf.ln(5)
    
    pdf.set_font("Arial", 'I', 10)
    pdf.multi_cell(0, 6, "Instructions: Read all questions carefully. Write your answers clearly.")
    pdf.ln(5)
    
    pdf.set_draw_color(200, 200, 200)
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    pdf.ln(5)

    pdf.set_font("Arial", size=11)
    for index, item in enumerate(export_questions, 1):
        pdf.set_font("Arial", 'B', 11)
        pdf.multi_cell(0, 6, remove_special_characters(f"{index}. {item['question']}"))
        pdf.set_font("Arial", size=11)
        for option_index, option in enumerate(item.get("options") or []):
            pdf.multi_cell(0, 6, remove_special_characters(f"   {chr(97 + option_index)}) {option}"))
        pdf.ln(3)

    if include_answers:
        pdf.add_page()
        pdf_heading(pdf, "Teacher Answer Sheet")
        for index, item in enumerate(export_questions, 1):
            pdf.set_font("Arial", 'B', 11)
            pdf.multi_cell(0, 6, remove_special_characters(f"{index}. {item['question']}"))
            pdf.set_font("Arial", size=10)
            pdf_label_value(pdf, "Correct Answer:", item.get("answer") or "Not provided")
            pdf_label_value(pdf, "Explanation:", item.get("explanation") or "Not provided")
            pdf_label_value(pdf, "Difficulty:", item.get("difficulty") or "Medium")
            pdf_label_value(pdf, "Source Reference:", (item.get("source_reference") or "Not provided")[:900])
            pdf.ln(4)

    pdf.add_page()
    pdf_heading(pdf, "AI Assessment Summary")
    pdf_label_value(pdf, "Total Questions:", summary["total_questions"])
    pdf_label_value(pdf, "Average Relevance:", summary["average_relevance"])
    pdf_label_value(pdf, "Average Coverage:", summary["average_coverage"])
    pdf_label_value(pdf, "Average Confidence:", summary["average_confidence"])
    difficulty_text = ", ".join(f"{key}: {value}" for key, value in summary["difficulty_distribution"].items())
    pdf_label_value(pdf, "Difficulty Distribution:", difficulty_text)
    pdf_label_value(pdf, "Topic Coverage Summary:", summary["topic_coverage_summary"])

    pdf_bytes = BytesIO()
    pdf_string = pdf.output(dest='S').encode('latin1')
    pdf_bytes.write(pdf_string)
    pdf_bytes.seek(0)
    return pdf_bytes


def save_questions_to_docx(raw_questions, subject, marks, include_answers=True):
    """Save questions to Word document (editable)"""
    doc = Document()
    
    # Title
    title = doc.add_heading('ASSESSMENT PAPER', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Info
    doc.add_paragraph(f"Subject: {subject}")
    doc.add_paragraph(f"Maximum Marks: {marks}")
    doc.add_paragraph()
    
    instructions = doc.add_paragraph("Instructions: Read all questions carefully. Write your answers clearly.")
    instructions.italic = True
    doc.add_paragraph("_" * 70)
    doc.add_paragraph()
    
    # Parse questions and answers
    question_pattern = r'(\d+\.)'
    parts = re.split(question_pattern, raw_questions)
    
    questions_only = []
    answers_only = []
    
    i = 1
    while i < len(parts):
        if re.match(r'\d+\.', parts[i]):
            q_num = parts[i]
            q_content = parts[i + 1] if i + 1 < len(parts) else ""
            question_text = q_num + q_content
            answer_text = ""
            
            for pattern in [r'\*\*Answer:.*', r'\*\*Solution:.*', r'\*\*Model Answer:.*', r'\*\*Answer Guidelines:.*']:
                match = re.search(pattern, question_text, re.IGNORECASE | re.DOTALL)
                if match:
                    answer_text = match.group(0).replace('**', '').strip()
                    question_text = question_text[:match.start()].strip()
                    break
            
            if question_text.strip():
                questions_only.append(question_text.strip())
                answers_only.append(answer_text.strip() if answer_text else "")
            i += 2
        else:
            i += 1
    
    # Add questions
    doc.add_heading('Questions', level=1)
    for q in questions_only:
        clean_q = re.sub(r'\*\*Answer:.*', '', q, flags=re.IGNORECASE)
        clean_q = re.sub(r'\*\*Solution:.*', '', clean_q, flags=re.IGNORECASE)
        clean_q = re.sub(r'\*\*Model Answer:.*', '', clean_q, flags=re.IGNORECASE)
        clean_q = re.sub(r'\*\*Answer Guidelines:.*', '', clean_q, flags=re.IGNORECASE)
        if clean_q.strip():
            doc.add_paragraph(clean_q.strip())
            doc.add_paragraph()
    
    # Add answers if requested
    if include_answers == True:
        doc.add_page_break()
        answer_title = doc.add_heading('TEACHER ANSWER SHEET', 0)
        answer_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph()
        
        for i, answer in enumerate(answers_only, 1):
            if answer.strip():
                ans_para = doc.add_paragraph()
                run = ans_para.add_run(f"Question {i}: ")
                run.bold = True
                ans_para.add_run(answer)
                doc.add_paragraph()
    
    docx_bytes = BytesIO()
    doc.save(docx_bytes)
    docx_bytes.seek(0)
    return docx_bytes

def remove_special_characters(text):
    """Remove special characters for PDF"""
    return unicodedata.normalize('NFKD', text).encode('ascii', 'ignore').decode('utf-8')

def create_google_form(questions, user_email):
    """Create Google Form (questions only, no answers)"""
    if not GOOGLE_FORMS_ENABLED:
        raise RuntimeError(f"Google Forms integration is not configured: {GOOGLE_FORMS_STATUS}")
        
    try:
        new_form = {
            "info": {
                "title": "Generated Assessment Form",
                "documentTitle": "Assessment Form"
            }
        }
        form = form_service.forms().create(body=new_form).execute()
        form_id = form['formId']
        form_url = form["responderUri"]

        if user_email:
            share_google_form(form_id, user_email)

        questions_list = parse_questions(questions)
        for q_data in questions_list:
            add_question_to_form(form_service, form_id, q_data['question'], q_data['type'], q_data.get('options'))

        return form_url
    except Exception as e:
        print(f"Error creating form: {e}")
        raise

def parse_questions(questions_text):
    """Parse questions (strip answers for forms)"""
    questions = []
    lines = questions_text.strip().split('\n')
    current_question = None
    current_options = []
    
    for line in lines:
        line = line.strip()
        if not line or '**Answer' in line or '**Solution' in line or '**Model' in line:
            continue
            
        if re.match(r'^\d+\.', line):
            if current_question:
                q_type = "mcq" if current_options else "text"
                questions.append({
                    'question': current_question,
                    'type': q_type,
                    'options': current_options if current_options else None
                })
            
            current_question = re.sub(r'^\d+\.\s*', '', line)
            current_options = []
            
        elif re.match(r'^[a-d]\)', line):
            option_text = re.sub(r'^[a-d]\)\s*', '', line)
            current_options.append(option_text)
    
    if current_question:
        q_type = "mcq" if current_options else "text"
        questions.append({
            'question': current_question,
            'type': q_type,
            'options': current_options if current_options else None
        })
    
    return questions

def share_google_form(form_id, user_email):
    """Share form with user"""
    if not GOOGLE_FORMS_ENABLED:
        return
        
    permission = {
        'type': 'user',
        'role': 'writer',
        'emailAddress': user_email
    }
    try:
        drive_service.permissions().create(fileId=form_id, body=permission, sendNotificationEmail=True).execute()
    except Exception as e:
        print(f"Error sharing form: {e}")

def add_question_to_form(service, form_id, question_text, question_type="text", choices=None):
    """Add question to form"""
    if not question_text.strip():
        return
    
    form = service.forms().get(formId=form_id).execute()
    existing_questions = form.get('items', [])
    question_count = len(existing_questions)

    if question_type == "mcq" and choices:
        question = {
            "requests": [{
                "createItem": {
                    "item": {
                        "title": question_text,
                        "questionItem": {
                            "question": {
                                "required": True,
                                "choiceQuestion": {
                                    "type": "RADIO",
                                    "options": [{"value": choice} for choice in choices]
                                }
                            }
                        }
                    },
                    "location": {"index": question_count}
                }
            }]
        }
    else:
        question = {
            "requests": [{
                "createItem": {
                    "item": {
                        "title": question_text,
                        "questionItem": {
                            "question": {
                                "required": True,
                                "textQuestion": {"paragraph": True}
                            }
                        }
                    },
                    "location": {"index": question_count}
                }
            }]
        }
    
    try:
        service.forms().batchUpdate(formId=form_id, body=question).execute()
    except Exception as e:
        print(f"Error adding question: {e}")

@app.route('/')
def index():
    return render_template('index.html', google_forms_enabled=GOOGLE_FORMS_ENABLED)

@app.route('/generate', methods=['POST'])
@rate_limit
def generate():
    try:
        uploaded_file = request.files.get('pdf_file')
        if not uploaded_file or uploaded_file.filename == '':
            return "No file uploaded", 400

        filename = uploaded_file.filename
        ext = filename.lower().split('.')[-1]
        
        # NEW: Check for supported file types
        if ext not in ['pdf', 'docx', 'pptx', 'txt']:
            return "Unsupported file format. Please upload PDF, DOCX, PPTX, or TXT files.", 400

        question_type = request.form.get('question_type', 'mcq')
        num_questions = int(request.form.get('num_questions', 10))
        output_format = request.form.get('output_format', 'pdf')
        subject = request.form.get('subject', 'General')
        marks = request.form.get('marks', '100')
        user_email = request.form.get('email', '').strip()
        include_answers = request.form.get('include_answers', 'false').lower() == 'true'
        difficulty = request.form.get('difficulty', 'Medium')
        chunk_size = int(request.form.get('chunk_size', os.getenv("RAG_CHUNK_SIZE", 500)))
        chunk_overlap = int(request.form.get('chunk_overlap', os.getenv("RAG_CHUNK_OVERLAP", 80)))
        top_k = min(
    int(request.form.get('top_k', os.getenv("RAG_TOP_K", 6))),
    3
)
        # Validation
        if num_questions < 1 or num_questions > 50:
            return "Number of questions must be between 1 and 50", 400

        if output_format == 'form':
            if not GOOGLE_FORMS_ENABLED:
                return "Google Forms integration is not configured", 400
            if not user_email:
                return "Email address is required for Google Form", 400
            if not validate_email(user_email):
                return "Invalid email format", 400

        # Save upload to the platform temp directory for serverless compatibility.
        temp_path = os.path.join(tempfile.gettempdir(), f"queryfy_{uuid.uuid4().hex}_{os.path.basename(filename)}")
        uploaded_file.save(temp_path)

        try:
            # NEW: Use the multi-format extraction function
            file_text = extract_text_from_file(temp_path, filename)
            
            if len(file_text.strip()) < 100:
                return "File content too short or unreadable. Please upload a file with at least 100 characters of text content.", 400

            doc_id, chunks, retrieved_chunks, context_chunks = build_rag_pipeline(
                file_text=file_text,
                filename=filename,
                subject=subject,
                question_type=question_type,
                difficulty=difficulty,
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                top_k=top_k,
            )
            db.store_document(
                doc_id=doc_id,
                filename=filename,
                file_type=ext,
                text_length=len(file_text),
                chunk_count=len(chunks),
                metadata={
                    "subject": subject,
                    "question_type": question_type,
                    "difficulty": difficulty,
                    "chunk_size": chunk_size,
                    "chunk_overlap": chunk_overlap,
                    "top_k": top_k,
                },
            )

            question_generator = QuestionGenerator(
                embedding_service=embedding_service,
                evaluation_service=evaluation_service,
                llm_client=llm_generate,
            )
            structured_questions = question_generator.generate(
                context_chunks=context_chunks,
                question_type=map_question_type(question_type),
                num_questions=num_questions,
                difficulty=difficulty,
            )

            if len(structured_questions) != num_questions:
                return (
                    f"Unable to generate exactly {num_questions} unique questions. "
                    f"Generated {len(structured_questions)} after regeneration attempts. "
                    "Please try a document with more content or reduce the requested count.",
                    500,
                )

            questions = convert_questions_to_text(structured_questions)

            global LAST_GENERATED_RESULTS
            LAST_GENERATED_RESULTS = []
            for item in structured_questions:
                metrics = {
                    "relevance_score": item.get("relevance_score", 0),
                    "context_coverage": item.get("context_coverage", 0),
                    "difficulty_score": item.get("difficulty_score", 0),
                    "confidence_score": item.get("confidence_score", 0),
                }
                db.store_question(
                    document_id=doc_id,
                    question_type=question_type,
                    difficulty=item.get("difficulty", difficulty),
                    question_text=item.get("question", ""),
                    answer=item.get("answer", ""),
                    explanation=item.get("explanation", ""),
                    source_chunk=item.get("source_excerpt", ""),
                    metrics=metrics,
                )
                LAST_GENERATED_RESULTS.append({
                    "question": item.get("question", ""),
                    "answer": item.get("answer", ""),
                    "explanation": item.get("explanation", ""),
                    "difficulty": item.get("difficulty", difficulty),
                    "source_chunk": item.get("source_excerpt", ""),
                    "retrieved_chunks": retrieved_chunks[:3],
                    **metrics,
                })

            if output_format == 'pdf':
                export_count = len(parse_questions_for_export(questions, structured_questions))
                if export_count != num_questions:
                    return f"PDF validation failed: requested {num_questions}, prepared {export_count} questions.", 500
                pdf_output = save_questions_to_pdf(
                    questions,
                    subject,
                    marks,
                    include_answers,
                    generated_items=structured_questions,
                )
                fname = f'{subject}_Assessment_with_Answers.pdf' if include_answers else f'{subject}_Assessment.pdf'
                return send_file(pdf_output, download_name=fname, as_attachment=True, mimetype='application/pdf')
            
            # NEW: Word document output
            elif output_format == 'docx':
                export_count = len(parse_questions_for_export(questions, structured_questions))
                if export_count != num_questions:
                    return f"DOCX validation failed: requested {num_questions}, prepared {export_count} questions.", 500
                docx_output = save_questions_to_docx(questions, subject, marks, include_answers)
                fname = f'{subject}_Assessment_with_Answers.docx' if include_answers else f'{subject}_Assessment.docx'
                return send_file(docx_output, download_name=fname, as_attachment=True, 
                               mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
            
            elif output_format == 'form':
                try:
                    export_count = len(parse_questions_for_export(questions, structured_questions))
                    if export_count != num_questions:
                        return jsonify({'success': False, 'message': f'Google Form validation failed: requested {num_questions}, prepared {export_count} questions.'}), 500
                    form_url = create_google_form(questions, user_email)
                    return jsonify({'success': True, 'form_url': form_url, 'message': 'Form created successfully!'})
                except Exception as e:
                    print(f"Google Form error: {e}")
                    return jsonify({'success': False, 'message': f'Google Form error: {str(e)}'}), 500
            
            else:
                return "Invalid output format", 400
        
        finally:
            if os.path.exists(temp_path):
                os.remove(temp_path)
            
    except ValueError as e:
        return f"Invalid input: {str(e)}", 400
    except RuntimeError as e:
        print(f"Runtime error: {str(e)}")
        return str(e), 503
    except Exception as e:
        print(f"Error: {str(e)}")
        return f"An error occurred: {str(e)}", 500

@app.route('/health')
def health():
    """Health check endpoint"""
    return jsonify({
        'status': 'healthy',
        'google_forms_enabled': GOOGLE_FORMS_ENABLED,
        'google_forms_status': GOOGLE_FORMS_STATUS,
        'rag_enabled': True,
        'vector_chunks_loaded': vector_store.count(),
        'vector_store_backend': getattr(vector_store, "backend", "memory"),
        'apis_configured': {
            'groq': bool(GROQ_API_KEY)
        },
        'groq_model': GROQ_MODEL,
        'startup_checks': STARTUP_CHECKS
    })

@app.route('/analytics')
def analytics():
    """Analytics dashboard data for processed documents and generated questions."""
    return jsonify(db.get_analytics())

@app.route('/results')
def results():
    """Latest generated question quality report for the existing UI."""
    return jsonify({'results': LAST_GENERATED_RESULTS})

if __name__ == '__main__':
    port = int(os.getenv('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=False)
